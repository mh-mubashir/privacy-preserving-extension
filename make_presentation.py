"""
make_presentation.py — Final presentation generator.
Run: python make_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Palette ──────────────────────────────────────────────────────────────────
BG       = RGBColor(0x12, 0x12, 0x24)   # very dark navy
PANEL    = RGBColor(0x1C, 0x1C, 0x3A)   # card background
HDR      = RGBColor(0x0D, 0x2B, 0x55)   # header bar
ACCENT   = RGBColor(0xE8, 0x4A, 0x35)   # coral red line
GOLD     = RGBColor(0xF5, 0xA6, 0x23)   # amber/gold
BLUE     = RGBColor(0x21, 0x96, 0xF3)   # bright blue
GREEN    = RGBColor(0x4C, 0xAF, 0x50)
ORANGE   = RGBColor(0xFF, 0x98, 0x00)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MUTED    = RGBColor(0xA0, 0xB4, 0xC8)   # muted text

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_slide():
    sl = prs.slides.add_slide(blank)
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = BG
    return sl


def rect(sl, l, t, w, h, fill, line=False):
    from pptx.util import Inches
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = fill
        s.line.width = Pt(0)
    else:
        s.line.fill.background()
    return s


def label(sl, text, l, t, w, h,
          size=16, bold=False, colour=WHITE,
          align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = colour
    return tb


def header(sl, title, subtitle=None):
    """Standard slide header: dark bar + red underline."""
    rect(sl, 0, 0, 13.33, 1.05, HDR)
    rect(sl, 0, 1.05, 13.33, 0.055, ACCENT)
    label(sl, title, 0.45, 0.08, 12.4, 0.7,
          size=30, bold=True, colour=WHITE)
    if subtitle:
        label(sl, subtitle, 0.45, 0.73, 12.4, 0.32,
              size=12, colour=MUTED)


def picture(sl, path, l, t, w, h):
    if os.path.exists(path):
        sl.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))


def card(sl, l, t, w, h, accent_col=None):
    rect(sl, l, t, w, h, PANEL)
    if accent_col:
        rect(sl, l, t, 0.15, h, accent_col)
    return l, t, w, h


def bottom_note(sl, text, colour=MUTED):
    rect(sl, 0, 7.1, 13.33, 0.4, HDR)
    label(sl, text, 0.4, 7.14, 12.5, 0.34,
          size=11, colour=colour, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, BG)
# horizontal decorative bars
rect(sl, 0, 2.55, 13.33, 0.05, ACCENT)
rect(sl, 0, 5.0,  13.33, 0.05, ACCENT)
# left side accent block
rect(sl, 0, 0, 0.22, 7.5, BLUE)

label(sl, "PRIVACY PRESERVING", 0.55, 0.9,  12.2, 1.1,
      size=56, bold=True, colour=WHITE, align=PP_ALIGN.LEFT)
label(sl, "EDGE VISION",         0.55, 1.82, 12.2, 1.0,
      size=56, bold=True, colour=GOLD, align=PP_ALIGN.LEFT)

label(sl, "EECE 5698  ·  Spring 2026",
      0.55, 2.8, 12.2, 0.5, size=18, colour=MUTED)
label(sl, "Hamza Mubashir   ·   Malia Howe   ·   Sindhu SureshKumar",
      0.55, 3.35, 12.2, 0.55, size=22, bold=True, colour=WHITE)
label(sl, "Northeastern University",
      0.55, 3.95, 12.2, 0.45, size=16, colour=MUTED)

# tag at bottom
rect(sl, 0.55, 5.3, 4.5, 0.6, BLUE)
label(sl, "Mid-Checkpoint Presentation", 0.7, 5.35, 4.2, 0.5,
      size=15, bold=True, colour=WHITE)

# ════════════════════════════════════════════════════════════════════════════
# 2 — AGENDA
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "Agenda")
items = [
    ("1", "Problem Statement",               BLUE),
    ("2", "Dataset & System Architecture",   BLUE),
    ("3", "VAE Variants & Architecture",     GOLD),
    ("4", "Compute Cost & Edge Deployment",  GOLD),
    ("5", "Experimental Results & Charts",   GREEN),
    ("6", "Two-Phase Training Curve",        GREEN),
    ("7", "Challenges & Why They Happened",  ORANGE),
    ("8", "Future Work",                     ACCENT),
]
cols = 2
rows_per_col = 4
for idx, (num, text, col) in enumerate(items):
    c = idx // rows_per_col
    r = idx % rows_per_col
    lx = 0.5 + c * 6.5
    ty = 1.25 + r * 1.45
    rect(sl, lx, ty, 5.9, 1.2, PANEL)
    rect(sl, lx, ty, 0.55, 1.2, col)
    label(sl, num, lx+0.04, ty+0.28, 0.48, 0.6,
          size=22, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
    label(sl, text, lx+0.7, ty+0.35, 5.1, 0.55,
          size=17, bold=True, colour=WHITE)

# ════════════════════════════════════════════════════════════════════════════
# 3 — PROBLEM STATEMENT
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "Problem Statement",
       "Privacy-preserving representation learning for edge vision systems")

# left card — goal
rect(sl, 0.4, 1.2, 5.9, 5.4, PANEL)
rect(sl, 0.4, 1.2, 0.15, 5.4, BLUE)
label(sl, "THE GOAL", 0.7, 1.35, 5.4, 0.42,
      size=13, bold=True, colour=BLUE)
goal_lines = [
    "Detect smiles reliably  (utility ✓)",
    "Suppress gender prediction  (privacy ✓)",
    "Run on edge hardware efficiently",
]
for i, line in enumerate(goal_lines):
    rect(sl, 0.65, 1.88 + i*0.82, 5.5, 0.62, HDR)
    label(sl, "▸  " + line, 0.82, 1.96 + i*0.82, 5.2, 0.48,
          size=14, colour=WHITE)

label(sl, "Target metric:", 0.65, 4.38, 5.4, 0.38,
      size=13, bold=True, colour=GOLD)
label(sl, "Privacy Accuracy ≈ 50%  (random guessing)\nUtility Accuracy ≥ 85%",
      0.65, 4.72, 5.4, 0.65, size=13, colour=MUTED)

# right card — challenge
rect(sl, 6.85, 1.2, 6.1, 5.4, PANEL)
rect(sl, 6.85, 1.2, 0.15, 5.4, ACCENT)
label(sl, "THE CHALLENGE", 7.15, 1.35, 5.7, 0.42,
      size=13, bold=True, colour=ACCENT)
label(sl, "Smile and gender are correlated in face images.\n\n"
          "A model trained to detect smiles will also inadvertently encode gender cues — "
          "facial structure, lip shapes, and skin texture all carry gender information.",
      7.15, 1.88, 5.65, 1.8, size=13, colour=WHITE)

label(sl, "Our Solution", 7.15, 3.85, 5.7, 0.38,
      size=14, bold=True, colour=GOLD)
label(sl, "Adversarial Representation Learning (ARL):\nTrain an encoder + adversary jointly. "
          "The adversary tries to predict gender from the encoder output. "
          "The encoder is trained to fool it — while still detecting smiles.",
      7.15, 4.28, 5.65, 1.2, size=13, colour=MUTED)

# ════════════════════════════════════════════════════════════════════════════
# 4 — DATASET & ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "Dataset & System Architecture")

# left: dataset
rect(sl, 0.4, 1.2, 4.1, 5.5, PANEL)
rect(sl, 0.4, 1.2, 0.15, 5.5, GOLD)
label(sl, "CelebA Dataset", 0.72, 1.32, 3.6, 0.45,
      size=15, bold=True, colour=GOLD)
dataset_items = [
    "202,599 celebrity face images",
    "40 binary attribute labels",
    "Utility label: Smiling (attr #31)",
    "Privacy label: Male (attr #20)",
    "Train 162k · Val 19k · Test 19k",
    "Input resolution: 224 × 224 RGB",
]
for i, item in enumerate(dataset_items):
    label(sl, "·  " + item, 0.72, 1.9 + i*0.72, 3.65, 0.58,
          size=13, colour=WHITE)

# right: ARL pipeline drawn as boxes+arrows
rect(sl, 4.9, 1.2, 8.05, 5.5, PANEL)
rect(sl, 4.9, 1.2, 0.15, 5.5, BLUE)
label(sl, "ARL Training Pipeline", 5.22, 1.32, 7.5, 0.45,
      size=15, bold=True, colour=BLUE)

pipeline = [
    (5.1, 2.1, 2.0, 0.75, "Input Image", "224×224 face", HDR),
    (7.5, 2.1, 2.0, 0.75, "VAE Encoder", "Compress → z", RGBColor(0x1B, 0x5E, 0x20)),
    (10.0, 2.1, 2.6, 0.75, "Decoder", "Reconstruct image", HDR),
    (5.8, 3.4, 2.5, 0.9, "Utility Classifier", "Predict Smile\n(maximize acc)", RGBColor(0x0D, 0x47, 0xA1)),
    (8.9, 3.4, 2.5, 0.9, "Privacy Adversary", "Predict Gender\n(fool this)", RGBColor(0x6A, 0x1A, 0x1A)),
]
for lx, ty, w, h, title, sub, col in pipeline:
    rect(sl, lx, ty, w, h, col)
    label(sl, title, lx+0.08, ty+0.08, w-0.12, 0.38,
          size=11, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
    label(sl, sub, lx+0.08, ty+0.42, w-0.12, 0.38,
          size=10, colour=MUTED, align=PP_ALIGN.CENTER)

# arrows
for arrow_x, arrow_y in [(7.1, 2.38), (9.6, 2.38)]:
    label(sl, "→", arrow_x, arrow_y, 0.4, 0.35, size=18, bold=True, colour=GOLD)

rect(sl, 5.2, 4.55, 6.6, 0.55, RGBColor(0x0A, 0x22, 0x40))
label(sl, "ARL Loss  =  Utility Loss  −  λ × Adversary Loss",
      5.3, 4.62, 6.4, 0.42, size=14, bold=True, colour=GOLD, align=PP_ALIGN.CENTER)

label(sl, "Two-phase training: warm-up (utility only) → ARL (adversary activated)",
      5.1, 5.3, 7.4, 0.42, size=12, colour=MUTED, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# 5 — VAE VARIANTS (Member 1)
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "VAE Encoder Variants", "Member 1 — Sindhu SureshKumar")

vaes1 = [
    ("Vanilla VAE",
     "Conv encoder  3→32→64→128→256→512 → bottleneck (μ, σ) → mirrored deconv decoder",
     "MSE(recon, input) + 1·KL",
     "Baseline. Lossy compression discards fine-grained gender cues from the latent representation.",
     BLUE),
    ("Beta-VAE  (β = 4)",
     "Identical architecture to VanillaVAE — same channel schedule, same decoder",
     "MSE + 4·KL  (4× stronger KL regularisation)",
     "Stronger KL forces the latent distribution closer to Gaussian. Gender and smile features are pushed to separate dimensions.",
     GREEN),
    ("Residual VAE",
     "ResNet-style residual blocks in encoder with skip connections. VanillaVAE decoder.",
     "MSE + 1·KL  (same loss as VanillaVAE)",
     "Residual connections allow richer feature hierarchies. Achieved best NAG score in our experiments.",
     GOLD),
]
for i, (name, arch, loss, note, col) in enumerate(vaes1):
    ty = 1.2 + i * 2.0
    rect(sl, 0.4, ty, 12.5, 1.78, PANEL)
    rect(sl, 0.4, ty, 0.15, 1.78, col)
    label(sl, name,  0.72, ty+0.12, 4.5, 0.48, size=16, bold=True, colour=col)
    label(sl, "Arch:  " + arch, 0.72, ty+0.56, 11.9, 0.38, size=12, colour=WHITE)
    label(sl, "Loss:  " + loss, 0.72, ty+0.92, 11.9, 0.35, size=12, colour=MUTED)
    label(sl, "▸  " + note,     0.72, ty+1.26, 11.9, 0.42, size=12, colour=GOLD, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# 6 — VAE VARIANTS (Members 2 & 3)
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "VAE Encoder Variants", "Member 2 — Hamza Mubashir  |  Member 3 — Malia Howe")

vaes2 = [
    ("Beta-TC VAE",
     "Wider channels  3→64→128→256→512→1024  with residual blocks throughout",
     "MSE + β_KL·MI(z;x) + β_TC·TC(z) + β_dim·dim-KL",
     "Explicit Total Correlation penalty forces independent latent dims — adversary can't use smile dims to infer gender.",
     RGBColor(0xAB, 0x47, 0xBC)),
    ("Disentangled Beta-VAE",
     "Skip-connection encoder. Latent z split into z_utility and z_privacy halves. Only z_utility is decoded.",
     "Standard β-VAE loss. Reconstruction computed from z_utility only.",
     "Most explicit privacy bottleneck: z_privacy acts as a sink where gender info is routed and never decoded.",
     ORANGE),
    ("VQ-VAE  (Member 3 — Malia Howe)",
     "Encoder → Vector Quantiser (discrete codebook, 256 entries, dim=64) → Decoder",
     "MSE + codebook loss + commitment loss  (no KL term)",
     "Discrete bottleneck is a strong information compressor. Collapsed in our runs — all inputs mapped to same entry.",
     ACCENT),
]
for i, (name, arch, loss, note, col) in enumerate(vaes2):
    ty = 1.2 + i * 2.0
    rect(sl, 0.4, ty, 12.5, 1.78, PANEL)
    rect(sl, 0.4, ty, 0.15, 1.78, col)
    label(sl, name,  0.72, ty+0.12, 6.0, 0.48, size=16, bold=True, colour=col)
    label(sl, "Arch:  " + arch, 0.72, ty+0.56, 11.9, 0.38, size=12, colour=WHITE)
    label(sl, "Loss:  " + loss, 0.72, ty+0.92, 11.9, 0.35, size=12, colour=MUTED)
    label(sl, "▸  " + note,     0.72, ty+1.26, 11.9, 0.42, size=12, colour=GOLD, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# 7 — COMPUTE COST
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "Compute Cost & Edge Deployment",
       "Measured with ONNX Runtime (real CPU inference, single-threaded) and thop (MACs/params)")

col_w  = [3.0, 1.25, 1.25, 1.65, 1.5, 1.5, 1.55]
col_lx = [0.3]
for w in col_w[:-1]:
    col_lx.append(col_lx[-1] + w)

hdrs = ["Encoder", "MACs", "Params", "CPU Latency", "Cortex-A76", "Jetson Nano", "Edge Ready"]

rect(sl, 0.3, 1.18, 12.73, 0.52, HDR)
for i, (h, x) in enumerate(zip(hdrs, col_lx)):
    label(sl, h, x+0.06, 1.22, col_w[i]-0.1, 0.42,
          size=12, bold=True, colour=GOLD, align=PP_ALIGN.CENTER)

rows = [
    ("Vanilla VAE",         "1.99G", "23.65M", "20ms",  "60ms",   "161ms", "✓", GREEN),
    ("Beta-VAE (β=4)",      "1.99G", "23.65M", "19ms",  "57ms",   "152ms", "✓", GREEN),
    ("Residual VAE",        "2.60G", "26.97M", "28ms",  "85ms",   "226ms", "✓", GREEN),
    ("VQ-VAE",              "3.48G", "0.66M",  "50ms",  "150ms",  "400ms", "~", ORANGE),
    ("Beta-TC VAE",         "9.98G", "38.14M", "159ms", "477ms", "1664ms", "✗", ACCENT),
    ("Disentangled β-VAE", "32.14G", "17.78M", "509ms","1528ms", "5357ms", "✗", ACCENT),
]
for r, (name, macs, params, cpu, ca76, jn, dep, col) in enumerate(rows):
    ty = 1.73 + r * 0.78
    bg_fill = PANEL if r % 2 == 0 else RGBColor(0x15, 0x15, 0x2C)
    rect(sl, 0.3, ty, 12.73, 0.76, bg_fill)
    vals = [name, macs, params, cpu, ca76, jn, dep]
    for i, (v, x) in enumerate(zip(vals, col_lx)):
        c = (GOLD if i == 0 else (col if i == 6 else WHITE))
        b = (i == 0)
        label(sl, v, x+0.06, ty+0.19, col_w[i]-0.1, 0.38,
              size=12, bold=b, colour=c, align=PP_ALIGN.CENTER)

bottom_note(sl, "Target: < 100ms on Jetson Nano for real-time edge.  VanillaVAE & BetaVAE meet this target.")

# ════════════════════════════════════════════════════════════════════════════
# 8 — RESULTS TABLE
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "Quantitative Results",
       "CelebA test set · 2,000 samples · Utility = Smile · Privacy = Gender")

cw2 = [2.9, 1.35, 1.35, 1.2, 1.2, 1.25, 2.05]
cx2 = [0.3]
for w in cw2[:-1]:
    cx2.append(cx2[-1] + w)

rect(sl, 0.3, 1.18, 12.73, 0.52, HDR)
for i, (h, x) in enumerate(zip(["Model", "Utility %", "Privacy %",
                                  "U-AUC", "P-AUC", "NAG", "Privacy Status"], cx2)):
    label(sl, h, x+0.06, 1.22, cw2[i]-0.1, 0.42,
          size=12, bold=True, colour=GOLD, align=PP_ALIGN.CENTER)

results = [
    ("VanillaVAE (λ=1)",     "86.1", "88.4", "0.946", "0.949", "0.940", "Leaking",   ACCENT),
    ("BetaVAE (λ=1)",        "81.5", "85.5", "0.911", "0.913", "0.886", "Leaking",   ACCENT),
    ("ResidualVAE (λ=1)",    "85.9", "86.4", "0.942", "0.941", "0.988", "Leaking",   ACCENT),
    ("ResidualVAE (λ=2)",    "86.9", "87.6", "0.941", "0.941", "0.980", "Leaking",   ORANGE),
    ("CVAE (3 epochs)",      "56.0", "60.6", "0.714", "0.529", "0.564", "Partial",   ORANGE),
    ("FactorVAE (3 epochs)", "46.1", "64.3", "0.595", "0.752", "0.000", "Collapsed", ACCENT),
    ("VQ-VAE",               "46.1", "76.2", "0.512", "0.509", "0.000", "Collapsed", ACCENT),
]
for r, row in enumerate(results):
    name, u, p, ua, pa, nag, status, col = row
    ty = 1.73 + r * 0.72
    bg_fill = PANEL if r % 2 == 0 else RGBColor(0x15, 0x15, 0x2C)
    rect(sl, 0.3, ty, 12.73, 0.70, bg_fill)
    vals = [name, u+"%", p+"%", ua, pa, nag, status]
    for i, (v, x) in enumerate(zip(vals, cx2)):
        c = (GOLD if i == 0 else
             (GREEN if i == 1 and float(u) >= 80 else
              (ORANGE if i == 2 and float(p) < 75 else
               (col if i == 6 else WHITE))))
        label(sl, v, x+0.06, ty+0.17, cw2[i]-0.1, 0.38,
              size=12, bold=(i == 0), colour=c, align=PP_ALIGN.CENTER)

bottom_note(sl, "Goal: Utility > 85%  and  Privacy ≈ 50%.  Best NAG: ResidualVAE (0.988)")

# ════════════════════════════════════════════════════════════════════════════
# 9 — UTILITY vs PRIVACY CHART
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "Utility vs Privacy Accuracy — All Models")
picture(sl, "docs/charts/chart_utility_privacy.png", 1.0, 1.2, 11.3, 5.65)
bottom_note(sl,
    "Privacy accuracy should be ≈50% (random).  All models score 85–90% — the adversary can still predict gender well.")

# ════════════════════════════════════════════════════════════════════════════
# 10 — NAG
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "NAG Score — Privacy-Utility Balance",
       "Normalized Accuracy Gap = utility gain / privacy leak  |  Higher = better  |  0 = collapsed model")
picture(sl, "docs/charts/chart_nag.png", 1.5, 1.2, 10.3, 5.3)
bottom_note(sl,
    "ResidualVAE achieves the best NAG.  FactorVAE & VQ-VAE score 0 because utility dropped to random chance (collapsed).")

# ════════════════════════════════════════════════════════════════════════════
# 11 — AUC
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "AUC-ROC — Utility vs Privacy Discriminability",
       "AUC = 0.5 → random guessing (ideal for privacy)  |  AUC = 1.0 → perfect discrimination")
picture(sl, "docs/charts/chart_auc.png", 1.2, 1.2, 11.0, 5.35)
bottom_note(sl,
    "All VAEs still allow near-perfect gender discrimination (AUC ~0.95).  VQ-VAE privacy AUC ≈ 0.5 due to collapse, not privacy.")

# ════════════════════════════════════════════════════════════════════════════
# 12 — LAMBDA SWEEP
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "Lambda (λ) Sweep — Does Stronger Penalty Help?",
       "ARL loss = utility_loss − λ × adversary_loss  |  Higher λ = more privacy pressure on encoder")
picture(sl, "docs/charts/chart_lambda_sweep.png", 0.9, 1.2, 11.5, 5.1)
bottom_note(sl,
    "Increasing λ from 1→3 barely moves privacy accuracy.  The bottleneck is the adversary's strength, not the penalty weight.")

# ════════════════════════════════════════════════════════════════════════════
# 13 — TRAINING CURVE
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "Two-Phase Training Curve — ResidualVAE (λ=2.0)",
       "Phase 1: encoder learns utility (adversary OFF)  →  Phase 2: adversary activated")

picture(sl, "docs/residualvae_training_curve.png", 0.4, 1.2, 8.2, 5.5)

# observation cards on the right
obs = [
    ("Warmup (ep 1-3)",    "Adversary scores below random chance (39%) — encoder is free to learn utility first", BLUE),
    ("Epoch 4 jump",       "Utility accuracy jumps to 82% in one epoch — warmup pre-training worked", GREEN),
    ("Both curves rise",   "After ARL activates, both utility and privacy rise together — core challenge", ORANGE),
    ("Final epoch 13",     "Utility 88.3%, Privacy 90.4% — adversary matches the encoder throughout", ACCENT),
]
for i, (title, desc, col) in enumerate(obs):
    ty = 1.25 + i * 1.52
    rect(sl, 8.85, ty, 4.1, 1.38, PANEL)
    rect(sl, 8.85, ty, 0.14, 1.38, col)
    label(sl, title, 9.12, ty+0.1,  3.7, 0.42, size=13, bold=True, colour=col)
    label(sl, desc,  9.12, ty+0.52, 3.7, 0.76, size=11, colour=MUTED)

# ════════════════════════════════════════════════════════════════════════════
# 14 — RECONSTRUCTIONS (VanillaVAE)
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "VAE Reconstructions — VanillaVAE",
       "10 epochs · 20k samples · Each strip: input face images (top) → encoder → decoder output (bottom)")

rect(sl, 0.3, 1.18, 2.0, 0.5, HDR)
label(sl, "Epoch 1", 0.38, 1.22, 1.85, 0.42, size=14, bold=True, colour=GOLD)
picture(sl, "visuals/member1_vanillavae_10e_20k_epoch01.png",
        0.3, 1.72, 12.73, 2.42)

rect(sl, 0.3, 4.3, 2.2, 0.5, HDR)
label(sl, "Epoch 10", 0.38, 4.34, 2.05, 0.42, size=14, bold=True, colour=GREEN)
picture(sl, "visuals/member1_vanillavae_10e_20k_epoch10.png",
        0.3, 4.84, 12.73, 2.42)

bottom_note(sl, "Reconstruction quality improves over epochs — sharper details, better colour fidelity at epoch 10")

# ════════════════════════════════════════════════════════════════════════════
# 15 — RECONSTRUCTIONS (BetaVAE)
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "VAE Reconstructions — BetaVAE (β=4)",
       "Stronger KL regularisation → smoother latent → slightly blurrier output than VanillaVAE")

rect(sl, 0.3, 1.18, 2.0, 0.5, HDR)
label(sl, "Epoch 1", 0.38, 1.22, 1.85, 0.42, size=14, bold=True, colour=GOLD)
picture(sl, "visuals/member1_betavae_10e_20k_epoch01.png",
        0.3, 1.72, 12.73, 2.42)

rect(sl, 0.3, 4.3, 2.2, 0.5, HDR)
label(sl, "Epoch 10", 0.38, 4.34, 2.05, 0.42, size=14, bold=True, colour=GREEN)
picture(sl, "visuals/member1_betavae_10e_20k_epoch10.png",
        0.3, 4.84, 12.73, 2.42)

bottom_note(sl,
    "β=4 trades reconstruction sharpness for better-structured latent space — utility is 81.5% vs VanillaVAE 86.1%")

# ════════════════════════════════════════════════════════════════════════════
# 16 — CHALLENGES
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "Challenges & Why They Happened")

challenges = [
    ("Privacy accuracy won't drop below ~85%",
     "Smile and gender are correlated in CelebA — lipstick, jaw structure, skin texture. The encoder cannot fully decouple them in 13 epochs with a simple ARL objective.",
     ACCENT),
    ("Adversary keeps pace with the encoder",
     "ARL trains both simultaneously. As the encoder improves its representations, the adversary improves too. They converge to an equilibrium where privacy is not suppressed.",
     ORANGE),
    ("VQ-VAE codebook collapse",
     "All face images mapped to the same codebook entry — model output became a constant average face. Both utility and privacy dropped to random chance. Root cause: learning rate / codebook size mismatch.",
     ACCENT),
    ("Latent adversary with 5 steps made privacy worse",
     "Latent z contains more raw gender signal than a reconstructed image. With 5 adversary steps per encoder step, adversary loss saturated at 0.09 — leaving the encoder no gradient signal to suppress gender.",
     ORANGE),
]
for i, (title, desc, col) in enumerate(challenges):
    ty = 1.22 + i * 1.5
    rect(sl, 0.4, ty, 12.5, 1.35, PANEL)
    rect(sl, 0.4, ty, 0.15, 1.35, col)
    label(sl, "⚠  " + title, 0.72, ty+0.1,  11.8, 0.45, size=14, bold=True, colour=col)
    label(sl, desc,           0.72, ty+0.56, 11.8, 0.68, size=12, colour=MUTED)

# ════════════════════════════════════════════════════════════════════════════
# 17 — FUTURE WORK
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "Future Work")

future = [
    ("Focus on one model",          "Fine-tune ResidualVAE (best NAG = 0.988) with more epochs and the full 160k training set.", GREEN),
    ("Gradient Reversal Layer",     "Replace ARL objective with a GRL — directly reverses gradients through the adversary, more principled than indirect penalisation.", BLUE),
    ("Lower resolution (64 × 64)", "Less pixel information = less gender signal available. Malia is exploring this for VQ-VAE to prevent codebook collapse.", ORANGE),
    ("INT8 Quantisation",           "Post-training quantisation would reduce latency ~4×. VanillaVAE (20ms) could hit real-time at < 10ms on edge CPU.", GOLD),
]
for i, (title, desc, col) in enumerate(future):
    ty = 1.22 + i * 1.48
    rect(sl, 0.4, ty, 12.5, 1.32, PANEL)
    rect(sl, 0.4, ty, 0.15, 1.32, col)
    label(sl, title, 0.72, ty+0.12,  5.5, 0.45, size=15, bold=True, colour=col)
    label(sl, desc,  0.72, ty+0.58, 11.8, 0.64, size=13, colour=MUTED)

# ════════════════════════════════════════════════════════════════════════════
# 18 — Q&A  (LAST SLIDE)
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, BG)
rect(sl, 0, 0, 0.22, 7.5, BLUE)
rect(sl, 0, 2.62, 13.33, 0.05, ACCENT)
rect(sl, 0, 4.85, 13.33, 0.05, ACCENT)

label(sl, "Thank You", 0.55, 1.1, 12.3, 1.1,
      size=60, bold=True, colour=WHITE, align=PP_ALIGN.LEFT)
label(sl, "Questions & Discussion", 0.55, 2.1, 12.3, 0.75,
      size=30, colour=GOLD, align=PP_ALIGN.LEFT)

label(sl, "Code & results available at:",
      0.55, 3.1, 12.3, 0.45, size=15, colour=MUTED)
label(sl, "github.com/mh-mubashir/privacy-preserving-extension  ·  branch: member1-sindhu",
      0.55, 3.52, 12.3, 0.45, size=15, bold=True, colour=WHITE)

label(sl, "Hamza Mubashir   ·   Malia Howe   ·   Sindhu SureshKumar",
      0.55, 5.25, 12.3, 0.5, size=20, bold=True, colour=WHITE, align=PP_ALIGN.LEFT)
label(sl, "Northeastern University  ·  EECE 5698  ·  Spring 2026",
      0.55, 5.8, 12.3, 0.45, size=15, colour=MUTED, align=PP_ALIGN.LEFT)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "Privacy_Preserving_Edge_Vision_Final.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
